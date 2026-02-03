from fastapi import UploadFile
import json
import logging
from io import BytesIO
from typing import Dict, List, Any, Tuple

# Import libraries for processing
# pdfminer.six: Used for robust PDF text extraction.
from pdfminer.high_level import extract_text as extract_text_from_pdf
# python-docx: Standard library for reading .docx files.
from docx import Document
# python-pptx: Standard library for reading .pptx files.
from pptx import Presentation

logger = logging.getLogger(__name__)

async def process_file(file: UploadFile) -> Dict[str, Any]:
    """
    Extracts text from the uploaded file and structures it.
    Input: UploadFile (PDF, DOCX, PPTX)
    Output: JSON structure with chapters/headings and paragraphs.
    
    Returns structure:
    {
        "filename": str,
        "chapters": [
            {
                "title": str,
                "paragraphs": List[str],
                "raw_text": str
            },
            ...
        ]
    }
    """
    filename = file.filename.lower()
    content = await file.read()
    
    data = {
        "filename": file.filename,
        "chapters": []
    }

    try:
        if filename.endswith(".pdf"):
            # Model/Library: pdfminer.six
            # Robust extraction of text from PDF streams
            text = extract_text_pdf(content)
            # Basic heuristic: Splitting by double newlines to simulate paragraphs/sections
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            data["chapters"].append({
                "title": "Extracted Content", 
                "paragraphs": paragraphs,
                "raw_text": text
            })
            
        elif filename.endswith(".docx"):
            # Model/Library: python-docx
            # Extracts text from Word documents, preserving paragraph structure
            paragraphs, text = extract_text_docx(content)
            data["chapters"].append({
                "title": "Extracted Content", 
                "paragraphs": paragraphs,
                "raw_text": text
            })
            
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            # Model/Library: python-pptx
            # Extracts text from PowerPoint slides, treating each slide as a chapter
            slides = extract_text_pptx(content)
            data["chapters"] = slides
            
        else:
            # Fallback for plain text files
            text = content.decode("utf-8", errors="ignore")
            paragraphs = [line.strip() for line in text.splitlines() if line.strip()]
            data["chapters"].append({
                "title": "Raw Text", 
                "paragraphs": paragraphs,
                "raw_text": text
            })
            
    except Exception as e:
        logger.error(f"Error processing file {filename}: {e}")
        return {"error": str(e)}

    return data

def extract_text_pdf(content: bytes) -> str:
    """Extract text from PDF bytes using pdfminer."""
    with BytesIO(content) as stream:
        text = extract_text_from_pdf(stream)
    return text

def extract_text_docx(content: bytes) -> Tuple[List[str], str]:
    """Extract text from DOCX bytes using python-docx. Returns (paragraphs_list, full_text)."""
    with BytesIO(content) as stream:
        doc = Document(stream)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        text = "\n".join(paragraphs)
    return paragraphs, text

def extract_text_pptx(content: bytes) -> List[Dict[str, Any]]:
    """Extract text from PPTX bytes using python-pptx, structuring by slide."""
    slides_data = []
    with BytesIO(content) as stream:
        prs = Presentation(stream)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            title = f"Slide {i+1}"
            
            # Extract title if available
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text
            
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if text:
                            slide_text.append(text)
            
            slides_data.append({
                "title": title,
                "paragraphs": slide_text,
                "raw_text": "\n".join(slide_text)
            })
    return slides_data
